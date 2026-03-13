"""
ppt_builder.py
- 챕터/슬라이드 계층 구조
- 챕터 간지 페이지 생성
- 본문 좌상단 "01 챕터제목" 표시
- multi_section 레이아웃: 섹션별 카드 텍스트박스
- 폰트: S-Core Dream 7 / Montserrat ExtraBold (타이틀)
        S-Core Dream 5 / Montserrat Medium (본문)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import re
import math


# ─── 유틸리티 ───
def hex_to_rgb(hex_color):
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def is_korean(text):
    return bool(re.search(r'[가-힣]', text))


def get_title_font(text):
    return 'S-Core Dream 7' if is_korean(text) else 'Montserrat ExtraBold'


def get_body_font(text):
    return 'S-Core Dream 5' if is_korean(text) else 'Montserrat Medium'


def set_title_style(para, text, size, color):
    para.text = text
    para.font.size = size
    para.font.color.rgb = color
    para.font.name = get_title_font(text)
    para.font.bold = False
    para.font.italic = False


def set_body_style(para, text, size, color):
    para.text = text
    para.font.size = size
    para.font.color.rgb = color
    para.font.name = get_body_font(text)
    para.font.bold = False
    para.font.italic = False


def add_run(paragraph, text, size, color, font_name=None):
    """한 paragraph 안에 run을 추가 (폰트 혼합 가능)"""
    run = paragraph.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.name = font_name or get_body_font(text)
    run.font.bold = False
    run.font.italic = False
    return run


class PPTBuilder:
    def __init__(self, theme="ir_pro", custom_colors=None):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide_count = 0
        self.theme = theme

        themes = {
            "ir_pro": {
                "bg": "#FFFFFF", "title": "#222222", "text": "#444444",
                "accent": "#1B2A4A", "accent2": "#E8EDF3", "card_bg": "#F4F6F9",
                "subtitle": "#888888", "label_bg": "#1B2A4A", "label_text": "#FFFFFF",
                "chapter_num": "#1B2A4A", "chapter_title": "#444444",
                "section_label_bg": "#1B2A4A", "section_label_text": "#FFFFFF",
                "table_header_bg": "#1B2A4A", "table_header_text": "#FFFFFF",
                "table_row_bg1": "#F4F6F9", "table_row_bg2": "#FFFFFF",
                "table_border": "#D0D8E0",
                "cover_overlay": "#1B2A4A", "divider_line": "#1B2A4A",
                "page_num": "#AAAAAA", "bullet_color": "#1B2A4A"
            },
            "dark_modern": {
                "bg": "#1A1A2E", "title": "#FFFFFF", "text": "#E0E0E0",
                "accent": "#E94560", "accent2": "#0F3460", "card_bg": "#16213E",
                "subtitle": "#B0B0B0", "label_bg": "#E94560", "label_text": "#FFFFFF",
                "chapter_num": "#E94560", "chapter_title": "#FFFFFF",
                "section_label_bg": "#E94560", "section_label_text": "#FFFFFF",
                "table_header_bg": "#E94560", "table_header_text": "#FFFFFF",
                "table_row_bg1": "#16213E", "table_row_bg2": "#1A1A2E",
                "table_border": "#0F3460",
                "cover_overlay": "#1A1A2E", "divider_line": "#E94560",
                "page_num": "#666666", "bullet_color": "#E94560"
            },
            "light_clean": {
                "bg": "#FFFFFF", "title": "#2D2D2D", "text": "#4A4A4A",
                "accent": "#4A90D9", "accent2": "#F0F4F8", "card_bg": "#F7F9FC",
                "subtitle": "#7A7A7A", "label_bg": "#4A90D9", "label_text": "#FFFFFF",
                "chapter_num": "#4A90D9", "chapter_title": "#2D2D2D",
                "section_label_bg": "#4A90D9", "section_label_text": "#FFFFFF",
                "table_header_bg": "#4A90D9", "table_header_text": "#FFFFFF",
                "table_row_bg1": "#F7F9FC", "table_row_bg2": "#FFFFFF",
                "table_border": "#D0D8E0",
                "cover_overlay": "#2D2D2D", "divider_line": "#4A90D9",
                "page_num": "#AAAAAA", "bullet_color": "#4A90D9"
            },
            "corporate_blue": {
                "bg": "#F0F4F8", "title": "#1A365D", "text": "#2D4A7A",
                "accent": "#3182CE", "accent2": "#EBF4FF", "card_bg": "#FFFFFF",
                "subtitle": "#4A6FA5", "label_bg": "#3182CE", "label_text": "#FFFFFF",
                "chapter_num": "#3182CE", "chapter_title": "#1A365D",
                "section_label_bg": "#3182CE", "section_label_text": "#FFFFFF",
                "table_header_bg": "#3182CE", "table_header_text": "#FFFFFF",
                "table_row_bg1": "#EBF4FF", "table_row_bg2": "#F0F4F8",
                "table_border": "#B0C4DE",
                "cover_overlay": "#1A365D", "divider_line": "#3182CE",
                "page_num": "#AAAAAA", "bullet_color": "#3182CE"
            },
        }

        self.colors = custom_colors if custom_colors else themes.get(theme, themes["ir_pro"])

    # ─── 공통 헬퍼 ───
    def _set_slide_bg(self, slide, color=None):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(color or self.colors["bg"])

    def _add_page_number(self, slide):
        self.slide_count += 1
        txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        set_body_style(p, str(self.slide_count), Pt(10), hex_to_rgb(self.colors.get("page_num", "#AAAAAA")))
        p.alignment = PP_ALIGN.RIGHT

    def _add_bottom_line(self, slide):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.95),
            Inches(11.733), Inches(0.015)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.colors.get("divider_line", "#DDDDDD"))
        line.line.fill.background()

    def _add_chapter_header(self, slide, chapter_num, chapter_title, slide_title=""):
        """좌상단에 챕터번호(accent) + 챕터제목(무채색) + 슬라이드 제목"""
        self._set_slide_bg(slide)

        num_str = f"{chapter_num:02d}"
        display_title = chapter_title

        # 챕터 번호 (accent 색상)
        tb_num = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(0.6), Inches(0.45))
        tf_n = tb_num.text_frame
        p_n = tf_n.paragraphs[0]
        set_title_style(p_n, num_str, Pt(20), hex_to_rgb(self.colors["chapter_num"]))

        # 챕터 제목 (무채색)
        tb_ct = slide.shapes.add_textbox(Inches(1.5), Inches(0.35), Inches(5.0), Inches(0.45))
        tf_ct = tb_ct.text_frame
        p_ct = tf_ct.paragraphs[0]
        set_body_style(p_ct, display_title, Pt(13), hex_to_rgb(self.colors["chapter_title"]))

        # 구분선
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85), Inches(0.6), Inches(0.03)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = hex_to_rgb(self.colors["chapter_num"])
        sep.line.fill.background()

        # 슬라이드 제목 (큰 폰트)
        if slide_title:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(11.5), Inches(0.65))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            set_title_style(p, slide_title, Pt(26), hex_to_rgb(self.colors["title"]))

        self._add_bottom_line(slide)

    # ═══════════════════════════════════════
    # 표지 슬라이드
    # ═══════════════════════════════════════
    def add_title_slide(self, title, subtitle="", bg_image_path=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        if bg_image_path and os.path.exists(bg_image_path):
            slide.shapes.add_picture(
                bg_image_path, Inches(0), Inches(0),
                self.prs.slide_width, self.prs.slide_height
            )
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                Inches(6.5), self.prs.slide_height
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = hex_to_rgb(self.colors["cover_overlay"])
            overlay.line.fill.background()
        else:
            self._set_slide_bg(slide, self.colors["cover_overlay"])

        # 상단 라벨
        lbl = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(5), Inches(0.5))
        tf_l = lbl.text_frame
        p_l = tf_l.paragraphs[0]
        set_body_style(p_l, "PRESENTATION", Pt(14), RGBColor(180, 190, 210))

        # 장식선
        deco = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.8), Inches(3.5), Inches(0.05)
        )
        deco.fill.solid()
        deco.fill.fore_color.rgb = RGBColor(255, 255, 255)
        deco.line.fill.background()

        # 타이틀
        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(3.1), Inches(5.0), Inches(1.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        set_title_style(p, title, Pt(40), RGBColor(255, 255, 255))

        # 부제
        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.9), Inches(5.0), Inches(0.8))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            set_body_style(p2, subtitle, Pt(18), RGBColor(200, 210, 225))

        self._add_page_number(slide)

    # ═══════════════════════════════════════
    # INDEX(목차) 슬라이드
    # ═══════════════════════════════════════
    def add_index_slide(self, chapters, bg_image_path=None):
        """chapters: ["챕터1 제목", "챕터2 제목", ...]"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        if bg_image_path and os.path.exists(bg_image_path):
            slide.shapes.add_picture(
                bg_image_path, Inches(0), Inches(0),
                self.prs.slide_width, self.prs.slide_height
            )
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                Inches(6.0), self.prs.slide_height
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = hex_to_rgb(self.colors["cover_overlay"])
            overlay.line.fill.background()
        else:
            self._set_slide_bg(slide, self.colors["cover_overlay"])

        # 타이틀
        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.0), Inches(4.5), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        set_title_style(p, "TABLE OF CONTENTS", Pt(16), RGBColor(180, 190, 210))

        start_top = Inches(2.0)
        for idx, ch_title in enumerate(chapters):
            num = f"{idx + 1:02d}"
            y = start_top + Inches(idx * 0.7)

            tb_n = slide.shapes.add_textbox(Inches(1.2), y, Inches(0.6), Inches(0.5))
            p_n = tb_n.text_frame.paragraphs[0]
            set_title_style(p_n, num, Pt(20), hex_to_rgb(self.colors.get("divider_line", "#FFFFFF")))

            tb_t = slide.shapes.add_textbox(Inches(2.0), y, Inches(3.5), Inches(0.5))
            tb_t.text_frame.word_wrap = True
            p_t = tb_t.text_frame.paragraphs[0]
            set_body_style(p_t, ch_title, Pt(16), RGBColor(255, 255, 255))

        self._add_page_number(slide)

    # ═══════════════════════════════════════
    # 챕터 간지(타이틀) 페이지
    # ═══════════════════════════════════════
    def add_chapter_title_slide(self, chapter_num, chapter_title, bg_image_path=None):
        """챕터 시작 간지 페이지"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.colors["cover_overlay"])

        if bg_image_path and os.path.exists(bg_image_path):
            slide.shapes.add_picture(
                bg_image_path, Inches(5), Inches(0),
                Inches(8.333), self.prs.slide_height
            )

        num_str = f"{chapter_num:02d}"

        # 챕터 번호 (큰 사이즈)
        tb_num = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(3.0), Inches(1.2))
        p_num = tb_num.text_frame.paragraphs[0]
        set_title_style(p_num, num_str, Pt(72), hex_to_rgb(self.colors.get("divider_line", "#FFFFFF")))

        # 구분선
        deco = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.3), Inches(2.5), Inches(0.05)
        )
        deco.fill.solid()
        deco.fill.fore_color.rgb = RGBColor(255, 255, 255)
        deco.line.fill.background()

        # 챕터 제목
        tb_title = slide.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(5.0), Inches(1.5))
        tb_title.text_frame.word_wrap = True
        p_title = tb_title.text_frame.paragraphs[0]
        set_title_style(p_title, chapter_title, Pt(32), RGBColor(255, 255, 255))

        self._add_page_number(slide)

    # ═══════════════════════════════════════
    # Multi-Section 레이아웃 (핵심 신규)
    # ═══════════════════════════════════════
    def add_multi_section_slide(self, chapter_num, chapter_title, slide_title,
                                 sections, image_path=None):
        """
        sections: [{"label": "Environmental", "lines": [...]}, ...]
        각 섹션을 별도 카드(텍스트박스)로 배치
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_chapter_header(slide, chapter_num, chapter_title, slide_title)

        num_sections = len(sections)
        has_image = image_path and os.path.exists(image_path)

        # 콘텐츠 영역 계산
        content_top = Inches(1.85)
        content_bottom = Inches(6.8)
        content_height = content_bottom - content_top

        if has_image:
            content_right = Inches(8.3)
            img_left = Inches(8.6)
            img_width = Inches(4.0)
            img_height = content_height
        else:
            content_right = Inches(12.5)

        content_left = Inches(0.8)
        total_width = content_right - content_left

        # 카드 배치 전략
        gap = Inches(0.2)

        if num_sections <= 3:
            # 가로 배치
            card_w = (total_width - gap * (num_sections - 1)) / num_sections
            card_h = content_height

            for idx, sec in enumerate(sections):
                x = content_left + idx * (card_w + gap)
                y = content_top
                self._draw_section_card(slide, x, y, card_w, card_h, sec)

        elif num_sections <= 6:
            # 2행 배치
            cols = math.ceil(num_sections / 2)
            card_w = (total_width - gap * (cols - 1)) / cols
            card_h = (content_height - gap) / 2

            for idx, sec in enumerate(sections):
                col = idx % cols
                row = idx // cols
                x = content_left + col * (card_w + gap)
                y = content_top + row * (card_h + gap)
                self._draw_section_card(slide, x, y, card_w, card_h, sec)
        else:
            # 3행 배치
            cols = math.ceil(num_sections / 3)
            card_w = (total_width - gap * (cols - 1)) / cols
            card_h = (content_height - gap * 2) / 3

            for idx, sec in enumerate(sections):
                col = idx % cols
                row = idx // cols
                x = content_left + col * (card_w + gap)
                y = content_top + row * (card_h + gap)
                self._draw_section_card(slide, x, y, card_w, card_h, sec)

        # 이미지
        if has_image:
            try:
                slide.shapes.add_picture(
                    image_path, img_left, content_top,
                    img_width, img_height
                )
            except Exception:
                pass

        self._add_page_number(slide)

    def _draw_section_card(self, slide, x, y, w, h, section):
        """섹션 카드: 라벨 바 + 내용 텍스트박스"""
        label = section.get("label", "")
        lines = section.get("lines", [])

        # 카드 배경
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(self.colors.get("card_bg", "#F4F6F9"))
        card.line.color.rgb = hex_to_rgb(self.colors.get("accent2", "#E8EDF3"))
        card.line.width = Pt(1)

        inner_margin = Inches(0.25)
        text_x = int(x) + inner_margin
        text_w = int(w) - inner_margin * 2
        current_top = int(y) + inner_margin

        # 라벨 바
        if label:
            label_h = Inches(0.38)
            label_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(label_h)
            )
            label_bar.fill.solid()
            label_bar.fill.fore_color.rgb = hex_to_rgb(
                self.colors.get("section_label_bg", self.colors["accent"])
            )
            label_bar.line.fill.background()

            tb_label = slide.shapes.add_textbox(
                text_x, int(y) + Inches(0.05), text_w, label_h
            )
            p_label = tb_label.text_frame.paragraphs[0]
            set_title_style(
                p_label, label, Pt(12),
                hex_to_rgb(self.colors.get("section_label_text", "#FFFFFF"))
            )
            current_top = int(y) + int(label_h) + Inches(0.1)

        # 내용
        remaining_h = int(y) + int(h) - current_top - inner_margin
        if remaining_h < Inches(0.3):
            remaining_h = Inches(0.3)

        tb_content = slide.shapes.add_textbox(text_x, current_top, text_w, remaining_h)
        tf = tb_content.text_frame
        tf.word_wrap = True

        # 폰트 크기 자동 조절
        total_lines = len(lines)
        if total_lines <= 3:
            fs = Pt(13)
            spacing = Pt(8)
        elif total_lines <= 5:
            fs = Pt(11)
            spacing = Pt(6)
        elif total_lines <= 8:
            fs = Pt(10)
            spacing = Pt(4)
        else:
            fs = Pt(9)
            spacing = Pt(3)

        bullet_c = hex_to_rgb(self.colors.get("bullet_color", "#1B2A4A"))
        text_c = hex_to_rgb(self.colors["text"])

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            # 불릿/기호 정리
            clean = line.strip()
            prefix = ""
            if clean.startswith(('· ', '• ', '- ', '– ')):
                prefix = "■  "
                clean = clean[2:].strip()
            elif clean.startswith('*'):
                prefix = "    ▸  "
                clean = clean[1:].strip()
            else:
                prefix = "■  "

            # Run 기반 스타일링 (불릿 색상 + 본문 색상 분리)
            p.text = ""
            run_bullet = p.add_run()
            run_bullet.text = prefix
            run_bullet.font.size = fs
            run_bullet.font.color.rgb = bullet_c
            run_bullet.font.name = get_body_font(prefix)
            run_bullet.font.bold = False
            run_bullet.font.italic = False

            run_text = p.add_run()
            run_text.text = clean
            run_text.font.size = fs
            run_text.font.color.rgb = text_c
            run_text.font.name = get_body_font(clean)
            run_text.font.bold = False
            run_text.font.italic = False

            p.space_after = spacing

    # ═══════════════════════════════════════
    # Bullets 레이아웃
    # ═══════════════════════════════════════
    def add_bullets_slide(self, chapter_num, chapter_title, slide_title,
                          points, image_path=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_chapter_header(slide, chapter_num, chapter_title, slide_title)

        has_img = image_path and os.path.exists(image_path)
        cw = Inches(7.0) if has_img else Inches(11.5)

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), cw, Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        n = len(points) if isinstance(points, list) else 1
        if n <= 3:
            fs, sp = Pt(17), Pt(14)
        elif n <= 5:
            fs, sp = Pt(15), Pt(10)
        elif n <= 8:
            fs, sp = Pt(13), Pt(8)
        else:
            fs, sp = Pt(11), Pt(6)

        if isinstance(points, list):
            for i, pt in enumerate(points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                set_body_style(p, f"■  {pt}", fs, hex_to_rgb(self.colors["text"]))
                p.space_after = sp
        else:
            p = tf.paragraphs[0]
            set_body_style(p, str(points), fs, hex_to_rgb(self.colors["text"]))

        if has_img:
            try:
                slide.shapes.add_picture(image_path, Inches(8.5), Inches(2.0), Inches(4.3), Inches(4.3))
            except Exception:
                pass

        self._add_page_number(slide)

    # ═══════════════════════════════════════
    # Two Column 레이아웃
    # ═══════════════════════════════════════
    def add_two_column_slide(self, chapter_num, chapter_title, slide_title,
                             left_points, right_points, image_path=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_chapter_header(slide, chapter_num, chapter_title, slide_title)

        col_w = Inches(5.5)
        col_top = Inches(2.1)

        for pts, x_off in [(left_points, Inches(0.8)), (right_points, Inches(6.6))]:
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_off, col_top, col_w, Inches(4.3))
            card.fill.solid()
            card.fill.fore_color.rgb = hex_to_rgb(self.colors.get("card_bg", "#F4F6F9"))
            card.line.color.rgb = hex_to_rgb(self.colors.get("accent2", "#E8EDF3"))
            card.line.width = Pt(1)

            txBox = slide.shapes.add_textbox(x_off + Inches(0.3), col_top + Inches(0.3),
                                             col_w - Inches(0.6), Inches(3.7))
            tf = txBox.text_frame
            tf.word_wrap = True
            fs = Pt(13) if len(pts) <= 5 else Pt(11)
            for i, pt in enumerate(pts):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                set_body_style(p, f"■  {pt}", fs, hex_to_rgb(self.colors["text"]))
                p.space_after = Pt(7)

        self._add_page_number(slide)

    # ═══════════════════════════════════════
    # Key-Value 레이아웃
    # ═══════════════════════════════════════
    def add_key_value_slide(self, chapter_num, chapter_title, slide_title,
                            pairs, image_path=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_chapter_header(slide, chapter_num, chapter_title, slide_title)

        max_cols = 3
        card_area = Inches(12.0)
        card_w = (card_area - Inches(0.3) * (max_cols - 1)) / max_cols
        card_h = Inches(2.0)

        for idx, pair in enumerate(pairs[:6]):
            col = idx % max_cols
            row = idx // max_cols
            c_left = Inches(0.6) + col * (card_w + Inches(0.3))
            c_top = Inches(2.2) + row * (card_h + Inches(0.3))

            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(c_left), c_top, int(card_w), card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = hex_to_rgb(self.colors.get("card_bg", "#F4F6F9"))
            card.line.color.rgb = hex_to_rgb(self.colors.get("accent2", "#E8EDF3"))
            card.line.width = Pt(1)

            tb_k = slide.shapes.add_textbox(int(c_left) + Inches(0.2), c_top + Inches(0.2),
                                            int(card_w) - Inches(0.4), Inches(0.5))
            pk = tb_k.text_frame.paragraphs[0]
            set_title_style(pk, pair.get("key", ""), Pt(15), hex_to_rgb(self.colors["chapter_num"]))

            tb_v = slide.shapes.add_textbox(int(c_left) + Inches(0.2), c_top + Inches(0.8),
                                            int(card_w) - Inches(0.4), Inches(1.0))
            tb_v.text_frame.word_wrap = True
            pv = tb_v.text_frame.paragraphs[0]
            set_body_style(pv, pair.get("value", ""), Pt(12), hex_to_rgb(self.colors["text"]))

        self._add_page_number(slide)

    # ═══════════════════════════════════════
    # Highlight 레이아웃
    # ═══════════════════════════════════════
    def add_highlight_slide(self, chapter_num, chapter_title, slide_title,
                            main_text, support_points=None, image_path=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_chapter_header(slide, chapter_num, chapter_title, slide_title)

        hl = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.0)
        )
        hl.fill.solid()
        hl.fill.fore_color.rgb = hex_to_rgb(self.colors["accent"])
        hl.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10.7), Inches(1.5))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        set_title_style(p, main_text, Pt(20), RGBColor(255, 255, 255))
        p.alignment = PP_ALIGN.CENTER

        if support_points:
            tb_sp = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.0))
            tb_sp.text_frame.word_wrap = True
            for i, sp in enumerate(support_points):
                ps = tb_sp.text_frame.paragraphs[0] if i == 0 else tb_sp.text_frame.add_paragraph()
                set_body_style(ps, f"■  {sp}", Pt(14), hex_to_rgb(self.colors["text"]))
                ps.space_after = Pt(8)

        self._add_page_number(slide)

    # ═══════════════════════════════════════
    # 표(Table) 레이아웃
    # ═══════════════════════════════════════
    def add_table_slide(self, chapter_num, chapter_title, slide_title,
                        headers, rows, image_path=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_chapter_header(slide, chapter_num, chapter_title, slide_title)

        num_rows = len(rows) + 1
        num_cols = len(headers)
        table_w = Inches(12.0)
        table_h = min(Inches(4.5), Inches(0.5) * num_rows)

        ts = slide.shapes.add_table(num_rows, num_cols, Inches(0.6), Inches(2.1), table_w, table_h)
        table = ts.table

        first_col_w = Inches(1.8)
        other_w = (table_w - first_col_w) / (num_cols - 1) if num_cols > 1 else table_w
        table.columns[0].width = int(first_col_w)
        for c in range(1, num_cols):
            table.columns[c].width = int(other_w)

        h_fs = Pt(11) if num_cols <= 6 else Pt(9)
        c_fs = Pt(10) if num_cols <= 6 else Pt(8)

        hdr_bg = hex_to_rgb(self.colors["table_header_bg"])
        hdr_tc = hex_to_rgb(self.colors["table_header_text"])
        r1_bg = hex_to_rgb(self.colors["table_row_bg1"])
        r2_bg = hex_to_rgb(self.colors["table_row_bg2"])
        txt_c = hex_to_rgb(self.colors["text"])

        for ci, ht in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = hdr_bg
            p = cell.text_frame.paragraphs[0]
            set_title_style(p, ht, h_fs, hdr_tc)
            p.alignment = PP_ALIGN.CENTER
            cell.text_frame.word_wrap = True
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for ri, rd in enumerate(rows):
            bg = r1_bg if ri % 2 == 0 else r2_bg
            for ci in range(num_cols):
                cell = table.cell(ri + 1, ci)
                ct = rd[ci] if ci < len(rd) else ""
                cell.text = ""
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                p = cell.text_frame.paragraphs[0]
                set_body_style(p, ct, c_fs, txt_c)
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                cell.text_frame.word_wrap = True
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        self._set_table_borders(table, num_rows, num_cols,
                                hex_to_rgb(self.colors["table_border"]))
        self._add_page_number(slide)

    def _set_table_borders(self, table, num_rows, num_cols, border_color):
        from pptx.oxml.ns import qn
        bc = str(border_color).replace('#', '')
        for ri in range(num_rows):
            for ci in range(num_cols):
                tc = table.cell(ri, ci)._tc
                tcPr = tc.get_or_add_tcPr()
                for bn in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                    ln = tcPr.find(qn(bn))
                    if ln is None:
                        ln = tcPr.makeelement(qn(bn), {})
                        tcPr.append(ln)
                    ln.set('w', '12700')
                    sf = ln.find(qn('a:solidFill'))
                    if sf is None:
                        sf = ln.makeelement(qn('a:solidFill'), {})
                        ln.insert(0, sf)
                    sc = sf.find(qn('a:srgbClr'))
                    if sc is None:
                        sc = sf.makeelement(qn('a:srgbClr'), {})
                        sf.append(sc)
                    sc.set('val', bc)

    # ═══════════════════════════════════════
    # 통합 렌더 (챕터 정보 포함)
    # ═══════════════════════════════════════
    def add_formatted_slide(self, formatted_data, chapter_num=1, chapter_title="",
                            image_path=None):
        layout = formatted_data.get("layout", "bullets")
        slide_title = formatted_data.get("title", "")
        sections = formatted_data.get("sections", [])

        if layout == "multi_section":
            self.add_multi_section_slide(
                chapter_num, chapter_title, slide_title,
                sections, image_path
            )
            return

        if layout == "table":
            for sec in sections:
                if sec["type"] == "table":
                    self.add_table_slide(
                        chapter_num, chapter_title, slide_title,
                        sec["headers"], sec["rows"], image_path
                    )
                    return

        if layout == "key_value":
            for sec in sections:
                if sec["type"] == "key_value":
                    self.add_key_value_slide(
                        chapter_num, chapter_title, slide_title,
                        sec["pairs"], image_path
                    )
                    return

        if layout == "two_column":
            lp, rp = [], []
            for sec in sections:
                if sec.get("side") == "left":
                    lp = sec.get("points", [])
                elif sec.get("side") == "right":
                    rp = sec.get("points", [])
            self.add_two_column_slide(
                chapter_num, chapter_title, slide_title,
                lp, rp, image_path
            )
            return

        if layout == "highlight":
            main, sup = "", []
            for sec in sections:
                if sec["type"] == "highlight":
                    main = sec["main"]
                elif sec["type"] == "support":
                    sup = sec.get("points", [])
            self.add_highlight_slide(
                chapter_num, chapter_title, slide_title,
                main, sup, image_path
            )
            return

        # 기본: bullets
        pts = []
        for sec in sections:
            if sec.get("type") == "bullets":
                pts.extend(sec.get("points", []))
            elif sec.get("type") == "text":
                pts.append(sec.get("content", ""))
        self.add_bullets_slide(
            chapter_num, chapter_title, slide_title,
            pts, image_path
        )

    # ═══════════════════════════════════════
    # 엔딩 슬라이드
    # ═══════════════════════════════════════
    def add_ending_slide(self, text="Thank You", subtext="", bg_image_path=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        if bg_image_path and os.path.exists(bg_image_path):
            slide.shapes.add_picture(
                bg_image_path, Inches(0), Inches(0),
                self.prs.slide_width, self.prs.slide_height
            )
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                self.prs.slide_width, self.prs.slide_height
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = hex_to_rgb(self.colors["cover_overlay"])
            overlay.line.fill.background()
        else:
            self._set_slide_bg(slide, self.colors["cover_overlay"])

        txBox = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(9), Inches(1.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        set_title_style(p, text, Pt(48), RGBColor(255, 255, 255))
        p.alignment = PP_ALIGN.CENTER

        if subtext:
            txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9), Inches(0.8))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            set_body_style(p2, subtext, Pt(18), RGBColor(200, 210, 225))
            p2.alignment = PP_ALIGN.CENTER

        self._add_page_number(slide)

    # ═══════════════════════════════════════
    # 저장
    # ═══════════════════════════════════════
    def save(self, filepath):
        d = os.path.dirname(filepath)
        if d and not os.path.exists(d):
            os.makedirs(d)
        self.prs.save(filepath)
        return filepath
