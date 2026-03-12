"""
generator/layout_engine.py
레이아웃 타입별 python-pptx 슬라이드 렌더링
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import math


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_rounded_rect(slide, left, top, width, height, color, radius=0.06):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    return shape


def is_korean(text):
    ko = sum(1 for ch in text if '\uac00' <= ch <= '\ud7a3' or '\u3130' <= ch <= '\u318f')
    en = sum(1 for ch in text if ch.isascii() and ch.isalpha())
    return ko >= en


FONT_KO = "S-Core Dream"
FONT_EN = "Montserrat SemiBold"


def get_font(text):
    return FONT_KO if is_korean(text) else FONT_EN


class LayoutEngine:
    def __init__(self, prs, theme, corner_radius=0.06, spacing="medium"):
        self.prs = prs
        self.theme = theme  # {"primary": RGBColor, "secondary": RGBColor, "accent": RGBColor}
        self.radius = corner_radius
        self.spacing = spacing
        self.page_num = 0

    def _sp(self):
        m = {
            "tight": {"margin": 0.4, "gap": 0.25},
            "medium": {"margin": 0.6, "gap": 0.35},
            "wide": {"margin": 0.8, "gap": 0.45},
            "extra-wide": {"margin": 1.0, "gap": 0.55},
        }
        return m.get(self.spacing, m["medium"])

    def _add_text(self, slide, left, top, width, height, text,
                  font_size=16, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or RGBColor(0, 0, 0)
        p.alignment = align
        p.font.name = get_font(str(text))
        return txBox

    def _add_slide_header(self, slide, title, label=""):
        sp = self._sp()
        # 왼쪽 악센트 바
        add_rounded_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5),
                         self.theme["primary"], radius=0)

        if label:
            self._add_text(slide, Inches(sp["margin"]), Inches(0.3),
                           Inches(8), Inches(0.3), label,
                           font_size=9, color=RGBColor(0x99, 0x99, 0x99))

        self._add_text(slide, Inches(sp["margin"]), Inches(0.7),
                       Inches(11), Inches(0.8), title,
                       font_size=28, bold=True, color=RGBColor(0x11, 0x11, 0x11))

        add_rounded_rect(slide, Inches(sp["margin"]), Inches(1.55),
                         Inches(0.9), Inches(0.04),
                         self.theme["accent"], radius=0.5)

    def _add_page_number(self, slide, num):
        txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(num)
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.name = FONT_EN

    def _new_slide(self):
        self.page_num += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        return slide

    # ──────────────────────────────────────
    # Bento Grid
    # ──────────────────────────────────────
    def render_bento_grid(self, data, label=""):
        slide = self._new_slide()
        self._add_slide_header(slide, data.get("title", ""), label)
        items = data.get("items", [])
        sp = self._sp()
        start_y = Inches(1.9)
        num = len(items)

        colors_cycle = [self.theme["primary"], self.theme["secondary"], self.theme["accent"]]

        if num <= 3:
            big_w = Inches(6.0)
            big_h = Inches(4.8)
            small_w = Inches(5.5)
            small_h = Inches(2.2)
            gap = Inches(sp["gap"])

            if items:
                add_rounded_rect(slide, Inches(sp["margin"]), start_y,
                                 big_w, big_h, colors_cycle[0], self.radius)
                self._add_text(slide, Inches(sp["margin"] + 0.4), start_y + Inches(0.4),
                               big_w - Inches(0.8), Inches(0.5), items[0].get("title", ""),
                               font_size=20, bold=True, color=RGBColor(255, 255, 255))
                self._add_text(slide, Inches(sp["margin"] + 0.4), start_y + Inches(1.1),
                               big_w - Inches(0.8), Inches(3.2), items[0].get("desc", ""),
                               font_size=14, color=RGBColor(220, 220, 220))

            right_x = Inches(sp["margin"]) + big_w + gap
            for i, item in enumerate(items[1:3]):
                y = start_y + i * (small_h + gap)
                c = colors_cycle[(i + 1) % len(colors_cycle)]
                add_rounded_rect(slide, right_x, y, small_w, small_h, c, self.radius)
                self._add_text(slide, right_x + Inches(0.3), y + Inches(0.25),
                               small_w - Inches(0.6), Inches(0.4), item.get("title", ""),
                               font_size=17, bold=True, color=RGBColor(255, 255, 255))
                self._add_text(slide, right_x + Inches(0.3), y + Inches(0.75),
                               small_w - Inches(0.6), Inches(1.2), item.get("desc", ""),
                               font_size=12, color=RGBColor(210, 210, 210))
        else:
            cols = min(num, 4)
            rows = math.ceil(num / cols)
            total_w = 12.333 - sp["margin"] * 2
            gap_w = sp["gap"]
            card_w = Inches((total_w - gap_w * (cols - 1)) / cols)
            card_h = Inches((5.0 - sp["gap"] * (rows - 1)) / rows)

            for i, item in enumerate(items[:8]):
                row, col = divmod(i, cols)
                x = Inches(sp["margin"]) + col * (card_w + Inches(gap_w))
                y = start_y + row * (card_h + Inches(sp["gap"]))
                c = colors_cycle[i % len(colors_cycle)]
                add_rounded_rect(slide, x, y, card_w, card_h, c, self.radius)
                self._add_text(slide, x + Inches(0.25), y + Inches(0.25),
                               card_w - Inches(0.5), Inches(0.45), item.get("title", ""),
                               font_size=16, bold=True, color=RGBColor(255, 255, 255))
                self._add_text(slide, x + Inches(0.25), y + Inches(0.8),
                               card_w - Inches(0.5), card_h - Inches(1.0), item.get("desc", ""),
                               font_size=11, color=RGBColor(210, 210, 210))

        self._add_page_number(slide, self.page_num)
        return slide

    # ──────────────────────────────────────
    # Stat Cards
    # ──────────────────────────────────────
    def render_stat_cards(self, data, label=""):
        slide = self._new_slide()
        self._add_slide_header(slide, data.get("title", ""), label)
        stats = data.get("stats", [])
        sp = self._sp()

        num = min(len(stats), 4)
        if num == 0:
            self._add_page_number(slide, self.page_num)
            return slide

        total_w = 12.333 - sp["margin"] * 2
        gap_w = sp["gap"]
        card_w = Inches((total_w - gap_w * (num - 1)) / num)
        card_h = Inches(4.0)
        start_y = Inches(2.2)

        for i, stat in enumerate(stats[:4]):
            x = Inches(sp["margin"]) + i * (card_w + Inches(gap_w))
            add_rounded_rect(slide, x, start_y, card_w, card_h,
                             self.theme["primary"], self.radius)
            self._add_text(slide, x + Inches(0.3), start_y + Inches(0.6),
                           card_w - Inches(0.6), Inches(1.8),
                           stat.get("number", ""),
                           font_size=52, bold=True,
                           color=self.theme["accent"], align=PP_ALIGN.CENTER)
            self._add_text(slide, x + Inches(0.3), start_y + Inches(2.5),
                           card_w - Inches(0.6), Inches(1.2),
                           stat.get("label", ""),
                           font_size=15, color=RGBColor(200, 200, 200),
                           align=PP_ALIGN.CENTER)

        self._add_page_number(slide, self.page_num)
        return slide

    # ──────────────────────────────────────
    # Timeline
    # ──────────────────────────────────────
    def render_timeline(self, data, label=""):
        slide = self._new_slide()
        self._add_slide_header(slide, data.get("title", ""), label)
        steps = data.get("steps", [])
        sp = self._sp()

        if not steps:
            self._add_page_number(slide, self.page_num)
            return slide

        bar_y = Inches(3.9)
        bar_left = Inches(1.0)
        bar_w = Inches(11.0)
        add_rounded_rect(slide, bar_left, bar_y, bar_w, Inches(0.05),
                         self.theme["secondary"], radius=0.5)

        num = min(len(steps), 6)
        step_w = 11.0 / num

        for i, step in enumerate(steps[:6]):
            cx = 1.0 + step_w * i + step_w / 2

            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx - 0.15), bar_y - Inches(0.1),
                Inches(0.3), Inches(0.3)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.theme["accent"]
            dot.line.fill.background()

            self._add_text(slide, Inches(cx - 0.9), Inches(2.3),
                           Inches(1.8), Inches(0.4),
                           step.get("year", f"Step {i+1}"),
                           font_size=13, bold=True,
                           color=self.theme["accent"], align=PP_ALIGN.CENTER)
            self._add_text(slide, Inches(cx - 0.9), Inches(2.7),
                           Inches(1.8), Inches(0.9),
                           step.get("title", ""),
                           font_size=14, bold=True,
                           color=self.theme["primary"], align=PP_ALIGN.CENTER)
            self._add_text(slide, Inches(cx - 1.0), Inches(4.4),
                           Inches(2.0), Inches(2.5),
                           step.get("desc", ""),
                           font_size=11, color=RGBColor(100, 100, 100),
                           align=PP_ALIGN.CENTER)

        self._add_page_number(slide, self.page_num)
        return slide

    # ──────────────────────────────────────
    # Quote
    # ──────────────────────────────────────
    def render_quote(self, data, label=""):
        slide = self._new_slide()
        sp = self._sp()

        add_rounded_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5),
                         self.theme["primary"], radius=0)

        add_rounded_rect(slide, Inches(1.5), Inches(1.5),
                         Inches(10.3), Inches(4.5),
                         self.theme["primary"], self.radius)

        self._add_text(slide, Inches(2.2), Inches(1.7),
                       Inches(1.0), Inches(1.2), "\u201C",
                       font_size=72, bold=True, color=self.theme["accent"])

        self._add_text(slide, Inches(2.5), Inches(2.8),
                       Inches(8.5), Inches(2.0),
                       data.get("quote", ""),
                       font_size=24, color=RGBColor(255, 255, 255),
                       align=PP_ALIGN.CENTER)

        author = data.get("author", "")
        if author:
            self._add_text(slide, Inches(2.5), Inches(5.0),
                           Inches(8.5), Inches(0.5),
                           f"\u2014 {author}",
                           font_size=15, color=self.theme["accent"],
                           align=PP_ALIGN.CENTER)

        self._add_page_number(slide, self.page_num)
        return slide

    # ──────────────────────────────────────
    # Big Number
    # ──────────────────────────────────────
    def render_big_number(self, data, label=""):
        slide = self._new_slide()
        sp = self._sp()

        add_rounded_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5),
                         self.theme["primary"], radius=0)

        add_rounded_rect(slide, Inches(2.0), Inches(1.2),
                         Inches(9.3), Inches(5.2),
                         self.theme["primary"], self.radius)

        self._add_text(slide, Inches(2.5), Inches(1.6),
                       Inches(8.3), Inches(2.5),
                       data.get("number", ""),
                       font_size=88, bold=True,
                       color=self.theme["accent"], align=PP_ALIGN.CENTER)

        self._add_text(slide, Inches(2.5), Inches(3.8),
                       Inches(8.3), Inches(1.2),
                       data.get("description", ""),
                       font_size=22, color=RGBColor(255, 255, 255),
                       align=PP_ALIGN.CENTER)

        subtitle = data.get("subtitle", "")
        if subtitle:
            self._add_text(slide, Inches(2.5), Inches(5.2),
                           Inches(8.3), Inches(0.8),
                           subtitle,
                           font_size=13, color=RGBColor(170, 170, 170),
                           align=PP_ALIGN.CENTER)

        self._add_page_number(slide, self.page_num)
        return slide

    # ──────────────────────────────────────
    # Icon Grid (Bento 변형)
    # ──────────────────────────────────────
    def render_icon_grid(self, data, label=""):
        data_copy = dict(data)
        data_copy["type"] = "bento_grid"
        return self.render_bento_grid(data_copy, label)

    # ──────────────────────────────────────
    # 통합 렌더
    # ──────────────────────────────────────
    def render(self, layout_type, data, label=""):
        renderers = {
            "bento_grid": self.render_bento_grid,
            "stat_cards": self.render_stat_cards,
            "timeline": self.render_timeline,
            "quote_slide": self.render_quote,
            "big_number": self.render_big_number,
            "icon_grid": self.render_icon_grid,
        }
        fn = renderers.get(layout_type)
        if fn:
            return fn(data, label)
        return None
